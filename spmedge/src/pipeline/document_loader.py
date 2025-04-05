#!/usr/bin/env python3
"""
Enhanced Document Loader - Second stage of the SPM Edge processing pipeline.
- Extracts content from various file formats with structure preservation
- Prepares documents for RAG by chunking and metadata generation
- Auto-detects document types and formats
- Provides quality metrics and extraction assessment
- Outputs enhanced JSON format optimized for downstream AI processing
"""

import os
import sys
import json
import logging
import shutil
import argparse
import re
from pathlib import Path
from datetime import datetime
from typing import Dict, Any, List, Optional, Union, Tuple

# Fix the import path issue - use absolute imports
sys.path.insert(0, '/Users/toddlebaron/prizym/spmedge')

# Now import using the absolute path
from config.config import config
from src.pipeline.db_integration import DBManager
from src.pipeline.pipeline_processor import PipelineProcessor, PipelineStage

# Ensure logs directory exists
LOG_DIR = config.LOG_DIR
os.makedirs(LOG_DIR, exist_ok=True)

LOG_FILE = os.path.join(LOG_DIR, "document_loader.log")

# Configure logging - disable propagation to avoid duplicates
logger = logging.getLogger("document_loader")
logger.propagate = False

if not logger.handlers:
    logger.setLevel(logging.INFO)
    formatter = logging.Formatter("%(asctime)s - %(levelname)s - %(message)s")
    
    # Add console handler
    console_handler = logging.StreamHandler()
    console_handler.setFormatter(formatter)
    logger.addHandler(console_handler)
    
    # File handler
    file_handler = logging.FileHandler(LOG_FILE, encoding="utf-8")
    file_handler.setFormatter(formatter)
    logger.addHandler(file_handler)

# Initialize database manager and pipeline processor
db_manager = DBManager()
processor = PipelineProcessor(PipelineStage.LOAD)


# ============= File Content Extraction Functions =============

def extract_text_from_pdf(pdf_path: Path) -> Dict[str, Any]:
    """
    Extract text content from a PDF file, preserving structure and metadata.
    
    Returns:
        Dict with extracted content, page count, and metadata
    """
    try:
        from PyPDF2 import PdfReader
        
        # Initialize extraction results
        text_content = ""
        pages = []
        metadata = {}
        
        with open(pdf_path, 'rb') as file:
            reader = PdfReader(file)
            
            # Extract document info
            info = reader.metadata
            if info:
                metadata = {
                    "title": info.title or "",
                    "author": info.author or "",
                    "subject": info.subject or "",
                    "creator": info.creator or "",
                    "producer": info.producer or ""
                }
            
            # Check if likely converted from another format
            likely_source = "native"
            if info and info.producer:
                if "Excel" in info.producer or "Spreadsheet" in info.producer:
                    likely_source = "spreadsheet"
                elif "Word" in info.producer or "Office" in info.producer:
                    likely_source = "word"
                elif "PowerPoint" in info.producer:
                    likely_source = "presentation"
            
            # Extract text by page
            for i, page in enumerate(reader.pages):
                page_text = page.extract_text() or ""
                pages.append({
                    "page_num": i+1,
                    "text": page_text,
                    "size": {"width": page.mediabox.width, "height": page.mediabox.height}
                })
                text_content += page_text + "\n\n"
        
        # Check if OCR might be needed (heuristic)
        needs_ocr = False
        if len(text_content.strip()) < 100 and len(pages) > 0:
            needs_ocr = True
        
        return {
            "content": text_content.strip() or "[Empty PDF content]",
            "page_count": len(pages),
            "pages": pages,
            "metadata": metadata,
            "likely_converted_from": likely_source,
            "needs_ocr": needs_ocr,
            "extraction_method": "PyPDF2",
            "extraction_quality": 0.9 if not needs_ocr else 0.3
        }
    except Exception as e:
        logger.error(f"❌ PDF extraction error for {pdf_path.name}: {e}")
        return {
            "content": f"[PDF extraction error: {str(e)}]",
            "page_count": 0,
            "pages": [],
            "metadata": {},
            "extraction_error": str(e),
            "extraction_method": "PyPDF2-failed",
            "needs_ocr": False,
            "extraction_quality": 0.0
        }


def extract_text_from_docx(docx_path: Path) -> Dict[str, Any]:
    """
    Extract text content from a Word document with structure preservation.
    
    Returns:
        Dict with extracted content, headings, and metadata
    """
    try:
        import docx
        
        doc = docx.Document(docx_path)
        full_text = []
        headings = []
        paragraphs = []
        
        # Extract document properties
        metadata = {}
        try:
            core_props = doc.core_properties
            metadata = {
                "title": core_props.title or "",
                "author": core_props.author or "",
                "subject": core_props.subject or "",
                "created": str(core_props.created) if core_props.created else "",
                "modified": str(core_props.modified) if core_props.modified else "",
                "category": core_props.category or ""
            }
        except:
            pass
        
        # Process paragraphs and identify headings
        for i, para in enumerate(doc.paragraphs):
            if para.text.strip():
                paragraphs.append({
                    "index": i,
                    "text": para.text,
                    "style": para.style.name if para.style else "Normal"
                })
                
                # Check if paragraph is a heading
                if para.style and "Heading" in para.style.name:
                    headings.append({
                        "text": para.text,
                        "level": int(para.style.name.replace("Heading ", "")) if para.style.name != "Heading" else 1,
                        "index": i
                    })
                
                full_text.append(para.text)
        
        # Extract tables
        tables = []
        for i, table in enumerate(doc.tables):
            table_data = []
            for row in table.rows:
                row_data = [cell.text for cell in row.cells]
                table_data.append(row_data)
            
            tables.append({
                "index": i,
                "data": table_data
            })
        
        content = "\n\n".join(full_text)
        
        return {
            "content": content or "[Empty Word document]",
            "headings": headings,
            "paragraphs": paragraphs,
            "tables": tables,
            "metadata": metadata,
            "extraction_method": "python-docx",
            "extraction_quality": 0.95
        }
    except ImportError:
        logger.warning("⚠️ python-docx not installed. Using basic extraction.")
        return {
            "content": f"[DOCX file: {docx_path.name} - Install python-docx for full extraction]",
            "extraction_method": "not-available",
            "extraction_quality": 0.0
        }
    except Exception as e:
        logger.error(f"❌ DOCX extraction error for {docx_path.name}: {e}")
        return {
            "content": f"[DOCX extraction error: {str(e)}]",
            "extraction_error": str(e),
            "extraction_method": "python-docx-failed",
            "extraction_quality": 0.0
        }


def extract_text_from_xlsx(xlsx_path: Path) -> Dict[str, Any]:
    """
    Extract text content from an Excel spreadsheet with structure preservation.
    
    Returns:
        Dict with extracted content, tables, and metadata
    """
    try:
        import pandas as pd
        
        # Initialize result
        sheets_data = []
        all_text = []
        
        # Read all sheets
        excel_file = pd.ExcelFile(xlsx_path)
        sheets = excel_file.sheet_names
        
        for sheet in sheets:
            df = pd.read_excel(excel_file, sheet_name=sheet)
            
            # Convert DataFrame to text representation
            sheet_text = f"Sheet: {sheet}\n{df.to_string(index=False)}\n"
            all_text.append(sheet_text)
            
            # Store structured data
            try:
                sheet_data = {
                    "name": sheet,
                    "headers": df.columns.tolist(),
                    "rows": df.values.tolist(),
                    "shape": df.shape
                }
                sheets_data.append(sheet_data)
            except Exception as e:
                logger.warning(f"⚠️ Error processing sheet {sheet}: {e}")
                sheets_data.append({
                    "name": sheet,
                    "error": str(e)
                })
        
        return {
            "content": "\n\n".join(all_text) or "[Empty spreadsheet]",
            "sheets": sheets_data,
            "sheet_count": len(sheets),
            "extraction_method": "pandas",
            "extraction_quality": 0.9
        }
    except ImportError:
        logger.warning("⚠️ pandas not installed. Using basic extraction.")
        return {
            "content": f"[XLSX file: {xlsx_path.name} - Install pandas for full extraction]",
            "extraction_method": "not-available",
            "extraction_quality": 0.0
        }
    except Exception as e:
        logger.error(f"❌ XLSX extraction error for {xlsx_path.name}: {e}")
        return {
            "content": f"[XLSX extraction error: {str(e)}]",
            "extraction_error": str(e),
            "extraction_method": "pandas-failed",
            "extraction_quality": 0.0
        }


def extract_text_from_pptx(pptx_path: Path) -> Dict[str, Any]:
    """
    Extract text content from a PowerPoint presentation with structure preservation.
    
    Returns:
        Dict with extracted content, slides, and metadata
    """
    try:
        from pptx import Presentation
        
        # Initialize extraction results
        all_text_content = []
        slides_data = []
        
        # Load presentation
        presentation = Presentation(pptx_path)
        
        # Process each slide
        for i, slide in enumerate(presentation.slides):
            slide_text = []
            shapes_data = []
            
            # Add slide number
            slide_text.append(f"Slide {i+1}")
            
            # Get slide title if available
            title = ""
            if slide.shapes.title and slide.shapes.title.text:
                title = slide.shapes.title.text
                slide_text.append(f"Title: {title}")
            
            # Extract text from all shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text)
                    shapes_data.append({
                        "text": shape.text,
                        "type": shape.__class__.__name__
                    })
            
            # Store slide data
            slide_content = "\n".join(slide_text)
            all_text_content.append(slide_content)
            
            slides_data.append({
                "number": i+1,
                "title": title,
                "content": slide_content,
                "shapes": shapes_data
            })
        
        return {
            "content": "\n\n".join(all_text_content) or "[Empty PowerPoint presentation]",
            "slides": slides_data,
            "slide_count": len(slides_data),
            "extraction_method": "python-pptx",
            "extraction_quality": 0.9
        }
    except ImportError:
        logger.warning("⚠️ python-pptx not installed. Using basic extraction.")
        return {
            "content": f"[PowerPoint file: {pptx_path.name} - Install python-pptx for full extraction]",
            "extraction_method": "not-available",
            "extraction_quality": 0.0
        }
    except Exception as e:
        logger.error(f"❌ PowerPoint extraction error for {pptx_path.name}: {e}")
        return {
            "content": f"[PowerPoint extraction error: {str(e)}]",
            "extraction_error": str(e),
            "extraction_method": "python-pptx-failed",
            "extraction_quality": 0.0
        }


def detect_converted_presentation(pdf_path: Path) -> bool:
    """
    Detect if a PDF is actually a converted PowerPoint presentation.
    Uses heuristics like slide numbers, layout patterns, etc.
    
    Args:
        pdf_path: Path to PDF file
        
    Returns:
        True if PDF appears to be a converted presentation
    """
    try:
        from PyPDF2 import PdfReader
        
        with open(pdf_path, 'rb') as file:
            reader = PdfReader(file)
            
            # Check document info
            info = reader.metadata
            if info and info.producer and any(term in info.producer for term in ["PowerPoint", "Keynote", "Slides", "Presentation"]):
                return True
            
            # Check page dimensions (presentations often have wider aspect ratio)
            if len(reader.pages) > 0:
                first_page = reader.pages[0]
                width = first_page.mediabox.width
                height = first_page.mediabox.height
                aspect_ratio = width / height
                
                # Presentations typically have 4:3 or 16:9 aspect ratios
                if 1.3 <= aspect_ratio <= 1.8:
                    # Look for slide indicators in text
                    text_samples = [reader.pages[i].extract_text() for i in range(min(5, len(reader.pages)))]
                    slide_patterns = [
                        r"slide\s+\d+",
                        r"^\s*\d+\s*$",  # Slide numbers
                        r"agenda|overview|summary|key takeaways|thank you"  # Common presentation terms
                    ]
                    
                    for sample in text_samples:
                        if sample and any(re.search(pattern, sample.lower()) for pattern in slide_patterns):
                            return True
            
            # Check for multiple pages with similar layouts
            if len(reader.pages) >= 3:
                # For presentations, pages often have similar structures
                page_sizes = set()
                for i in range(min(5, len(reader.pages))):
                    page = reader.pages[i]
                    page_sizes.add((page.mediabox.width, page.mediabox.height))
                
                if len(page_sizes) <= 2:  # Most slides have consistent dimensions
                    return True
            
            return False
    except Exception as e:
        logger.warning(f"⚠️ Error detecting presentation format: {e}")
        return False


def detect_converted_spreadsheet(pdf_path: Path) -> bool:
    """
    Detect if a PDF is actually a converted Excel spreadsheet.
    Uses heuristics like grid patterns, cell references, etc.
    
    Args:
        pdf_path: Path to PDF file
        
    Returns:
        True if PDF appears to be a converted spreadsheet
    """
    try:
        from PyPDF2 import PdfReader
        
        with open(pdf_path, 'rb') as file:
            reader = PdfReader(file)
            
            # Check document info
            info = reader.metadata
            if info and info.producer and any(term in info.producer for term in ["Excel", "Spreadsheet", "Calc", "Numbers"]):
                return True
            
            # Extract text to look for spreadsheet indicators
            if len(reader.pages) > 0:
                # Sample first few pages
                text_samples = [reader.pages[i].extract_text() for i in range(min(3, len(reader.pages)))]
                
                # Look for spreadsheet patterns
                spreadsheet_patterns = [
                    r"[A-Z]+\d+",  # Cell references like A1, B12
                    r"=\s*[A-Z]+\d+",  # Formulas
                    r"sum\(|average\(|count\(",  # Common spreadsheet functions
                    r"\$[\d,.]+\s+\$[\d,.]+",  # Currency columns
                    r"\d+\.\d+%\s+\d+\.\d+%",  # Percentage columns
                    r"^\s*\d+\s+[A-Za-z]",  # Row numbers with text
                    r"total|subtotal|grand total"  # Common spreadsheet terms
                ]
                
                pattern_matches = 0
                for sample in text_samples:
                    if not sample:
                        continue
                    for pattern in spreadsheet_patterns:
                        if re.search(pattern, sample, re.IGNORECASE):
                            pattern_matches += 1
                            if pattern_matches >= 3:  # If we find multiple patterns
                                return True
                
                # Check for grid-like text layout
                for sample in text_samples:
                    if sample:
                        lines = sample.split('\n')
                        if len(lines) > 5:
                            # Count lines with consistent spacing/tabs
                            consistent_lines = 0
                            spaces = []
                            for line in lines:
                                # Look for regular patterns of spaces that would indicate columns
                                space_positions = [m.start() for m in re.finditer(r'\s{2,}', line)]
                                if space_positions:
                                    spaces.append(space_positions)
                            
                            # Check if space positions are consistent across lines
                            if spaces and len(spaces) > 5:
                                # Calculate similarity between consecutive space patterns
                                similar_lines = 0
                                for i in range(1, len(spaces)):
                                    # If lines have similar spacing, likely a table
                                    if len(spaces[i]) == len(spaces[i-1]):
                                        similar_lines += 1
                                
                                if similar_lines > 3:  # Several lines with similar structure
                                    return True
            
            return False
    except Exception as e:
        logger.warning(f"⚠️ Error detecting spreadsheet format: {e}")
        return False


def extract_from_presentation_pdf(pdf_path: Path) -> Dict[str, Any]:
    """
    Extract content from a PDF that was originally a presentation.
    Uses specialized extraction to preserve slide structure.
    
    Args:
        pdf_path: Path to PDF file
        
    Returns:
        Dict with extracted content optimized for presentations
    """
    # Start with standard PDF extraction
    result = extract_text_from_pdf(pdf_path)
    
    try:
        from PyPDF2 import PdfReader
        
        with open(pdf_path, 'rb') as file:
            reader = PdfReader(file)
            
            # Initialize slides data structure
            slides = []
            
            # Process each page as a slide
            for i, page in enumerate(reader.pages):
                page_text = page.extract_text() or ""
                
                # Try to identify slide title
                title = ""
                lines = page_text.split('\n')
                if lines and len(lines) > 0:
                    # First non-empty line is often the title in presentations
                    for line in lines:
                        if line.strip() and not line.strip().isdigit():  # Skip page numbers
                            title = line.strip()
                            break
                
                # Create slide object
                slide = {
                    "number": i+1,
                    "title": title,
                    "content": page_text,
                    "text_blocks": []
                }
                
                # Try to identify text blocks/bullet points
                bullet_pattern = r"^[\s•\-\*]+(.+)$"
                bullet_points = []
                current_section = ""
                
                for line in lines:
                    line = line.strip()
                    if not line:
                        continue
                        
                    # Skip the title we already identified
                    if line == title:
                        continue
                    
                    # Look for bullet points
                    bullet_match = re.match(bullet_pattern, line)
                    if bullet_match:
                        bullet_points.append(bullet_match.group(1).strip())
                    elif line.endswith(':'):  # Potential section header
                        current_section = line
                        if bullet_points:
                            slide["text_blocks"].append({
                                "type": "bullets",
                                "content": bullet_points.copy()
                            })
                            bullet_points = []
                    elif not re.match(r"^\d+$", line):  # Skip page numbers
                        # Regular text
                        if bullet_points:
                            slide["text_blocks"].append({
                                "type": "bullets",
                                "content": bullet_points.copy()
                            })
                            bullet_points = []
                        
                        slide["text_blocks"].append({
                            "type": "text",
                            "section": current_section,
                            "content": line
                        })
                
                # Add any remaining bullet points
                if bullet_points:
                    slide["text_blocks"].append({
                        "type": "bullets",
                        "content": bullet_points
                    })
                
                slides.append(slide)
            
            # Update result with presentation-specific data
            result["slides"] = slides
            result["slide_count"] = len(slides)
            result["detected_format"] = "converted_presentation"
            result["extraction_method"] = "presentation_pdf_extraction"
            
            # Rebuild content with better slide structure
            slide_texts = []
            for slide in slides:
                slide_text = [f"SLIDE {slide['number']}: {slide['title']}"]
                
                for block in slide["text_blocks"]:
                    if block["type"] == "text":
                        slide_text.append(block["content"])
                    elif block["type"] == "bullets":
                        for bullet in block["content"]:
                            slide_text.append(f"• {bullet}")
                
                slide_texts.append("\n".join(slide_text))
            
            result["content"] = "\n\n".join(slide_texts)
            result["extraction_quality"] = 0.85  # Better than regular PDF extraction
            
    except Exception as e:
        logger.error(f"❌ Error in presentation PDF extraction: {e}")
        # Fall back to standard extraction result
        
    return result


def extract_from_spreadsheet_pdf(pdf_path: Path) -> Dict[str, Any]:
    """
    Extract content from a PDF that was originally a spreadsheet.
    Uses specialized extraction to preserve table structure.
    
    Args:
        pdf_path: Path to PDF file
        
    Returns:
        Dict with extracted content optimized for spreadsheets
    """
    # Start with standard PDF extraction
    result = extract_text_from_pdf(pdf_path)
    
    try:
        from PyPDF2 import PdfReader
        import tabula
        import pandas as pd
        
        # Use tabula-py to extract tables
        tables = tabula.read_pdf(str(pdf_path), pages='all', multiple_tables=True)
        
        extracted_tables = []
        for i, table in enumerate(tables):
            if not table.empty:
                # Convert DataFrame to structured format
                extracted_tables.append({
                    "table_id": i+1,
                    "headers": table.columns.tolist(),
                    "data": table.values.tolist(),
                    "shape": table.shape
                })
        
        # If tables were extracted successfully
        if extracted_tables:
            result["tables"] = extracted_tables
            result["table_count"] = len(extracted_tables)
            result["detected_format"] = "converted_spreadsheet"
            result["extraction_method"] = "spreadsheet_pdf_extraction"
            
            # Rebuild content with better table formatting
            content_parts = []
            with open(pdf_path, 'rb') as file:
                reader = PdfReader(file)
                
                for i, page in enumerate(reader.pages):
                    page_text = page.extract_text() or ""
                    
                    # Try to extract page title/header
                    lines = page_text.split('\n')
                    page_title = ""
                    if lines and len(lines) > 0:
                        page_title = lines[0].strip()
                    
                    content_parts.append(f"PAGE {i+1}: {page_title}")
                    
                    # Add tables found on this page
                    page_tables = [t for t in extracted_tables if 
                                  # This is a simple heuristic - tables from this page
                                  # For more accuracy, we'd need to map tables to pages
                                  t["table_id"] <= (i+1)*3 and t["table_id"] > i*3]
                    
                    if page_tables:
                        for table in page_tables:
                            content_parts.append(f"TABLE {table['table_id']}:")
                            table_text = []
                            # Add headers
                            table_text.append(" | ".join(str(h) for h in table["headers"]))
                            table_text.append("-" * 50)  # Separator line
                            
                            # Add data rows
                            for row in table["data"]:
                                table_text.append(" | ".join(str(cell) for cell in row))
                            
                            content_parts.append("\n".join(table_text))
                    else:
                        # Use original page text if no tables detected
                        content_parts.append(page_text)
            
            # Update content with better structure
            result["content"] = "\n\n".join(content_parts)
            result["extraction_quality"] = 0.85  # Better than regular PDF extraction
    except ImportError:
        logger.warning("⚠️ tabula-py not installed. Using basic spreadsheet extraction.")
    except Exception as e:
        logger.error(f"❌ Error in spreadsheet PDF extraction: {e}")
        # Fall back to standard extraction result
        
    return result


def detect_and_process_document(file_path: Path, metadata: Dict[str, Any] = None) -> Dict[str, Any]:
    """
    Detect document format and use appropriate extraction method.
    Works with actual format rather than just file extension.
    
    Args:
        file_path: Path to document file
        metadata: Optional metadata from previous processing
        
    Returns:
        Dict with extracted content and enhanced metadata
    """
    if metadata is None:
        metadata = {}
    
    file_ext = file_path.suffix.lower().lstrip('.')
    metadata["file_type"] = file_ext
    
    # First detect the actual format regardless of extension
    detected_format = None
    
    if file_ext == "pdf":
        # Check if PDF is actually a converted presentation
        if detect_converted_presentation(file_path):
            detected_format = "converted_presentation"
            logger.info(f"📊 Detected converted presentation: {file_path.name}")
        # Check if PDF is actually a converted spreadsheet
        elif detect_converted_spreadsheet(file_path):
            detected_format = "converted_spreadsheet"
            logger.info(f"📈 Detected converted spreadsheet: {file_path.name}")
    
    # Update metadata with detected format
    if detected_format:
        metadata["detected_format"] = detected_format
    
    # Process based on actual format rather than extension
    if detected_format == "converted_presentation":
        result = extract_from_presentation_pdf(file_path)
    elif detected_format == "converted_spreadsheet":
        result = extract_from_spreadsheet_pdf(file_path)
    elif file_ext == "pdf":
        result = extract_text_from_pdf(file_path)
    elif file_ext == "docx":
        result = extract_text_from_docx(file_path)
    elif file_ext in ["xlsx", "xls"]:
        result = extract_text_from_xlsx(file_path)
    elif file_ext in ["pptx", "ppt"]:
        result = extract_text_from_pptx(file_path)
    elif file_ext in ["txt", "md", "csv"]:
        # Simple text files
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()
            result = {
                "content": content,
                "extraction_method": "direct_read",
                "extraction_quality": 1.0
            }
    elif file_ext == "json":
        # JSON files
        with open(file_path, 'r', encoding='utf-8') as f:
            json_content = json.load(f)
            result = {
                "content": json.dumps(json_content),
                "structured_data": json_content,
                "extraction_method": "json_parse",
                "extraction_quality": 1.0
            }
    else:
        # Unsupported format
        result = {
            "content": f"[Unsupported file format: {file_ext}]",
            "extraction_method": "unsupported",
            "extraction_quality": 0.0
        }
    
    # Add basic file info
    result["filename"] = file_path.name
    result["file_type"] = file_ext
    result["file_size"] = file_path.stat().st_size
    result["extraction_time"] = datetime.now().isoformat()
    
    # Calculate word count
    if "content" in result:
        result["word_count"] = len(re.findall(r'\w+', result["content"]))
    
    # Add detected format if available
    if detected_format:
        result["detected_format"] = detected_format
    
    # Set extraction success flag
    result["extraction_success"] = result.get("extraction_quality", 0) > 0 and "content" in result
    
    return result


# ============= RAG Preparation Functions =============

def chunk_document(content: str, chunk_size: int = 1000, overlap: int = 200) -> List[Dict[str, Any]]:
    """
    Split document into overlapping chunks for RAG preparation.
    
    Args:
        content: Document content to split
        chunk_size: Target size of each chunk in characters
        overlap: Number of characters to overlap between chunks
        
    Returns:
        List of chunks with text and position info
    """
    chunks = []
    
    # Simple paragraph-based chunking
    paragraphs = content.split("\n\n")
    
    current_chunk = []
    current_size = 0
    chunk_id = 1
    
    for para in paragraphs:
        para = para.strip()
        if not para:
            continue
            
        para_size = len(para)
        
        # If adding this paragraph exceeds chunk size and we already have content
        if current_size + para_size > chunk_size and current_chunk:
            # Create a chunk from what we have
            chunk_text = "\n\n".join(current_chunk)
            chunks.append({
                "chunk_id": chunk_id,
                "text": chunk_text,
                "position": chunks[-1]["position"] + len(chunks[-1]["text"]) - overlap if chunks else 0,
                "size": len(chunk_text),
                "paragraph_count": len(current_chunk)
            })
            chunk_id += 1
            
            # Start a new chunk with overlap
            overlap_size = min(overlap, len(chunk_text))
            if overlap_size > 0:
                # Keep some of the previous paragraphs for context
                overlap_paras = []
                overlap_content_size = 0
                for prev_para in reversed(current_chunk):
                    if overlap_content_size + len(prev_para) <= overlap:
                        overlap_paras.insert(0, prev_para)
                        overlap_content_size += len(prev_para)
                    else:
                        break
                
                current_chunk = overlap_paras
                current_size = overlap_content_size
            else:
                current_chunk = []
                current_size = 0
        
        # Add the paragraph to the current chunk
        current_chunk.append(para)
        current_size += para_size
    
    # Don't forget the last chunk
    if current_chunk:
        chunk_text = "\n\n".join(current_chunk)
        chunks.append({
            "chunk_id": chunk_id,
            "text": chunk_text,
            "position": chunks[-1]["position"] + len(chunks[-1]["text"]) - overlap if chunks else 0,
            "size": len(chunk_text),
            "paragraph_count": len(current_chunk)
        })
    
    return chunks


def detect_document_type(content: str, filename: str) -> Dict[str, Any]:
    """
    Detect document type based on content and filename.
    Uses simple rule-based heuristics.
    
    Returns:
        Dict with document type info
    """
    filename = filename.lower()
    content_sample = content[:5000].lower() if content else ""
    
    # Check for compensation plan documents
    if "comp_plan" in filename or "compensation" in filename:
        confidence = 0.9
        return {
            "detected_type": "compensation_plan",
            "confidence": confidence,
            "detection_method": "filename_pattern"
        }
    elif "quota" in filename or "target" in filename:
        confidence = 0.85
        return {
            "detected_type": "quota_model",
            "confidence": confidence,
            "detection_method": "filename_pattern"
        }
    
    # Content-based detection
    if re.search(r"compensation plan|comp plan|incentive plan", content_sample):
        matches = len(re.findall(r"compensation|incentive|bonus|commission|payout", content_sample))
        confidence = min(0.6 + matches * 0.05, 0.9)
        return {
            "detected_type": "compensation_plan",
            "confidence": confidence,
            "detection_method": "content_pattern",
            "pattern_matches": matches
        }
    elif re.search(r"quota|sales target|revenue goal", content_sample):
        matches = len(re.findall(r"quota|target|goal|revenue|sales", content_sample))
        confidence = min(0.6 + matches * 0.05, 0.9)
        return {
            "detected_type": "quota_model",
            "confidence": confidence,
            "detection_method": "content_pattern",
            "pattern_matches": matches
        }
    elif re.search(r"sales rep|territory|account", content_sample):
        return {
            "detected_type": "sales_document",
            "confidence": 0.7,
            "detection_method": "content_pattern"
        }
    
    # Default to unknown type
    return {
        "detected_type": "unknown",
        "confidence": 0.5,
        "detection_method": "default"
    }


def prepare_rag_document(document_id: str, original_filename: str, 
                        extraction_result: Dict[str, Any], 
                        document_type: str = None) -> Dict[str, Any]:
    """
    Prepare a document for RAG by adding chunks and metadata.
    
    Args:
        document_id: Document identifier
        original_filename: Original filename
        extraction_result: Result from extraction process
        document_type: Document type if known
        
    Returns:
        Dict with RAG-ready document
    """
    content = extraction_result.get("content", "")
    
    # Auto-detect document type if not provided
    doc_type_info = {}
    if document_type and document_type != "unknown":
        doc_type_info = {
            "detected_type": document_type,
            "confidence": 1.0,
            "detection_method": "predefined"
        }
    else:
        doc_type_info = detect_document_type(content, original_filename)
    
    # Create chunks for RAG
    chunks = chunk_document(content)
    
    # Compute statistics
    stats = {
        "word_count": extraction_result.get("word_count", 0),
        "chunk_count": len(chunks),
        "extraction_quality": extraction_result.get("extraction_quality", 0.0),
    }
    
    # Build the RAG document
    rag_doc = {
        "document_id": document_id,
        "content": content,
        "metadata": {
            "original_filename": original_filename,
            "file_type": extraction_result.get("file_type", ""),
            "file_size": extraction_result.get("file_size", 0),
            "extraction_method": extraction_result.get("extraction_method", ""),
            "extraction_quality": extraction_result.get("extraction_quality", 0.0),
            "structure_preserved": extraction_result.get("structure_preserved", False)
        },
        "document_type": doc_type_info,
        "chunks": chunks,
        "stats": stats
    }
    
    # Include structured data if available
    if "pages" in extraction_result:
        rag_doc["structured_data"] = {
            "pages": extraction_result["pages"]
        }
    if "slides" in extraction_result:
        rag_doc["structured_data"] = {
            "slides": extraction_result["slides"]
        }
    if "sheets" in extraction_result:
        rag_doc["structured_data"] = {
            "sheets": extraction_result["sheets"]
        }
    if "headings" in extraction_result:
        rag_doc["structured_data"] = rag_doc.get("structured_data", {})
        rag_doc["structured_data"]["headings"] = extraction_result["headings"]
    if "tables" in extraction_result:
        rag_doc["structured_data"] = rag_doc.get("structured_data", {})
        rag_doc["structured_data"]["tables"] = extraction_result["tables"]
    
    return rag_doc

def get_batch_size_from_settings(db_manager, default_limit=500):
    """
    Retrieve the batch size from pipeline settings or use default.
    
    Args:
        db_manager: Database manager instance
        default_limit: Default limit to use if setting not found
        
    Returns:
        int: Batch size from settings or default
    """
    try:
        db_manager.cursor.execute(
            "SELECT value FROM pipeline_settings WHERE key = 'batch.size';"
        )
        result = db_manager.cursor.fetchone()
        
        if result and result[0]:
            try:
                batch_size = int(result[0])
                logger.info(f"Using batch size from settings: {batch_size}")
                return batch_size
            except ValueError:
                logger.warning(f"Invalid batch size in settings: {result[0]}, using default: {default_limit}")
                return default_limit
        else:
            logger.info(f"No batch size found in settings, using default: {default_limit}")
            return default_limit
    except Exception as e:
        logger.warning(f"Error retrieving batch size from settings: {e}, using default: {default_limit}")
        return default_limit
        
def get_document_type(document_id: str) -> Optional[str]:
    """Get document type from the database."""
    try:
        result = db_manager.get_document_type(document_id)
        return result
    except Exception as e:
        logger.error(f"❌ Error getting document type for {document_id}: {e}")
        return None

def load_documents(limit: int = 100, output_format: str = "json", retry_failed: bool = False) -> int:
    """
    Load and process documents that have completed the input stage.
    
    Args:
        limit: Maximum number of documents to process
        output_format: Output format (json is required for RAG)
        retry_failed: Whether to retry previously failed documents
        
    Returns:
        Number of successfully processed documents
    """
    try:
        # Get documents that completed the INPUT stage
        documents = []
        
        if retry_failed:
            # Get both completed and failed documents from the input stage
            completed_docs = processor.get_documents_for_stage(current_stage="input", status="completed", limit=limit)
            failed_docs = processor.get_documents_for_stage(current_stage="load", status="failed", limit=limit)
            
            # Combine both lists (giving priority to completed ones if limit is reached)
            documents = completed_docs + failed_docs
            if len(documents) > limit:
                documents = documents[:limit]
            
            logger.info(f"Found {len(completed_docs)} completed docs and {len(failed_docs)} failed docs (retry mode)")
        else:
            # Just get completed documents from input stage
            documents = processor.get_documents_for_stage(current_stage="input", status="completed", limit=limit)
        
        if not documents:
            logger.info("No documents found ready for loading.")
            return 0

        # Create batch for this processing run
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        batch_name = f"load_batch_{timestamp}"
        batch_id = processor.record_batch_processing(
            batch_name=batch_name, 
            document_count=len(documents), 
            status="processing"
        )

        dirs = processor.get_base_dirs()
        stage_input_dir = dirs["stage_input"]  # Stage-specific input directory
        stage_load_dir = dirs["stage_load"]    # Stage-specific load directory
        
        # Ensure directories exist
        stage_input_dir.mkdir(parents=True, exist_ok=True)
        stage_load_dir.mkdir(parents=True, exist_ok=True)
        
        # Track processed documents
        loaded_documents = []
        success_count = 0
        error_count = 0
        missing_count = 0

        # Process each document
        for doc in documents:
            try:
                document_id = doc['id']
                original_filename = doc['name']
                
                logger.info(f"Processing document {document_id}: {original_filename}")
                processor.update_document_stage(document_id, status="processing", batch_id=batch_id)

                # Get document type from database
                document_type = get_document_type(document_id) or "unknown"
                
                # Find the file in stage_input directory
                file_path = stage_input_dir / original_filename
                
                if not file_path.exists():
                    # If not found, check original input directory as fallback
                    fallback_path = dirs["input"] / original_filename
                    
                    if fallback_path.exists():
                        file_path = fallback_path
                        logger.info(f"File found in fallback location: {fallback_path}")
                    else:
                        # Check in unprocessed directory as a last resort
                        unprocessed_path = dirs["unprocessed"] / original_filename
                        if unprocessed_path.exists():
                            file_path = unprocessed_path
                            logger.info(f"File found in unprocessed directory: {unprocessed_path}")
                        else:
                            # Check in stage_load directory - for retry operations
                            load_path = dirs["stage_load"] / original_filename
                            if load_path.exists():
                                file_path = load_path
                                logger.info(f"File found in stage_load directory: {load_path}")
                            else:
                                # Try to find original file backup
                                original_backup = dirs["stage_load"] / f"original_{Path(original_filename).name.split('_')[-1]}"
                                if original_backup.exists():
                                    file_path = original_backup
                                    logger.info(f"Found original backup file: {original_backup}")
                                else:
                                    # Try to find any file that might match by document ID
                                    id_short = str(document_id).replace('-', '')[:12]
                                    potential_files = []
                                    for search_dir in [stage_input_dir, dirs["input"], dirs["unprocessed"], dirs["stage_load"]]:
                                        potential_files.extend(list(search_dir.glob(f"*{id_short}*")))
                                    
                                    if potential_files:
                                        file_path = potential_files[0]
                                        logger.info(f"Found potential matching file: {file_path}")
                                    else:
                                        logger.error(f"❌ File not found: {file_path} or any fallbacks")
                                        processor.update_document_stage(
                                            document_id, 
                                            status="failed", 
                                            error_message="File not found", 
                                            batch_id=batch_id
                                        )
                                        missing_count += 1
                                        continue

                # Use intelligent document detection and processing
                extraction_result = detect_and_process_document(file_path)
                
                # Check if extraction was successful
                if not extraction_result.get("extraction_success"):
                    logger.warning(f"⚠️ Content extraction issues for {file_path.name}: {extraction_result.get('extraction_message', 'Unknown error')}")
                    
                # Prepare document for RAG
                rag_document = prepare_rag_document(
                    document_id=document_id,
                    original_filename=original_filename,
                    extraction_result=extraction_result,
                    document_type=document_type
                )
                
                # Convert to desired output format
                output_content = json.dumps(rag_document, indent=2) if output_format == "json" else rag_document["content"]

                # Generate new filename for load stage
                new_filename = processor.generate_stage_filename(
                    original_filename=file_path.name, 
                    document_id=document_id, 
                    batch_id=batch_id
                )
                
                # Update file extension based on output format
                if output_format == "json" and not new_filename.endswith(".json"):
                    new_filename = Path(new_filename).stem + ".json"
                elif output_format == "markdown" and not new_filename.endswith((".md", ".markdown")):
                    new_filename = Path(new_filename).stem + ".md"
                
                # Create the new file in load stage directory
                new_file_path = stage_load_dir / new_filename
                
                # Write content to new file
                with open(new_file_path, 'w', encoding='utf-8') as f:
                    f.write(output_content)
                logger.info(f"✅ Created output file: {new_file_path}")
                
                # Copy original file to load stage as backup if different
                if file_path.name != new_filename:
                    backup_path = stage_load_dir / f"original_{file_path.name}"
                    shutil.copy(str(file_path), str(backup_path))
                    logger.info(f"📁 Copied original file to: {backup_path}")

                # Update document in database
                db_manager.cursor.execute(
                    """
                    UPDATE documents 
                    SET name = %s, 
                        updated_at = NOW()
                    WHERE id = %s;
                    """, 
                    (new_filename, document_id)
                )
                db_manager.conn.commit()

                # Create metadata for database
                metadata = {
                    "extraction_time": extraction_result.get("extraction_time"),
                    "file_type": extraction_result.get("file_type"),
                    "file_size": extraction_result.get("file_size"),
                    "detected_format": extraction_result.get("detected_format", ""),
                    "extraction_method": extraction_result.get("extraction_method"),
                    "extraction_quality": extraction_result.get("extraction_quality", 0),
                    "content_format": output_format,
                    "document_type": document_type,
                    "detected_document_type": rag_document["document_type"]["detected_type"],
                    "detection_confidence": rag_document["document_type"]["confidence"],
                    "original_filename": original_filename,
                    "pipeline_filename": new_filename,
                    "word_count": extraction_result.get("word_count", 0),
                    "chunk_count": len(rag_document.get("chunks", [])),
                    "structure_preserved": extraction_result.get("structure_preserved", False),
                    "needs_ocr": extraction_result.get("needs_ocr", False)
                }
                
                # Save metadata to database
                db_manager.cursor.execute(
                    """
                    UPDATE documents 
                    SET metadata = %s
                    WHERE id = %s;
                    """,
                    (json.dumps(metadata), document_id)
                )
                db_manager.conn.commit()

                # Add to processed list for batch file
                loaded_documents.append({
                    "id": document_id,
                    "name": new_filename,
                    "original_filename": original_filename,
                    "document_type": document_type,
                    "detected_document_type": rag_document["document_type"]["detected_type"],
                    "detection_confidence": rag_document["document_type"]["confidence"],
                    "content": rag_document["content"][:1000] + ("..." if len(rag_document["content"]) > 1000 else ""),  # Truncated preview
                    "metadata": metadata,
                    "stats": rag_document["stats"],
                    "batch_id": batch_id,
                    "pipeline_stage": "load",
                    "status": "completed"
                })

                # Update document status in pipeline
                processor.update_document_stage(document_id, status="completed", batch_id=batch_id)
                logger.info(f"✅ Successfully processed document: {document_id}")
                success_count += 1

            except Exception as e:
                logger.error(f"❌ Error processing document {doc.get('id')}: {e}")
                processor.update_document_stage(
                    doc.get('id'), 
                    status="failed", 
                    error_message=str(e), 
                    batch_id=batch_id
                )
                error_count += 1

        # Save batch results
        if loaded_documents:
            processor.save_document_batch(documents=loaded_documents, batch_name=batch_name)
            logger.info(f"✅ Successfully loaded {success_count} out of {len(documents)} documents.")
            
            # Update batch status
            if success_count == len(documents):
                processor.finalize_batch(batch_id, "completed")
            else:
                processor.finalize_batch(batch_id, "partial")
                
            # Log error summary
            if error_count > 0 or missing_count > 0:
                logger.warning(f"⚠️ Document loading summary: {success_count} succeeded, {error_count} failed with errors, {missing_count} missing files")
        else:
            logger.warning("⚠️ No documents were successfully loaded.")
            processor.finalize_batch(batch_id, "failed")
        
        return success_count

    except Exception as e:
        logger.error(f"❌ Failed to load documents: {e}")
        return 0

def main():
    """Main function to execute document loading."""
    parser = argparse.ArgumentParser(description="Enhanced Document Loader - Extracts content and prepares for RAG")
    parser.add_argument("--limit", "-l", type=int, default=None, 
                        help="Maximum number of documents to load (overrides DB setting)")
    parser.add_argument("--format", "-f", type=str, default="json", 
                        choices=["text", "json", "markdown", "md"],
                        help="Output format for document content")
    parser.add_argument("--retry", "-r", action="store_true",
                        help="Retry previously failed documents") 
    parser.add_argument("--verbose", "-v", action="store_true",
                        help="Enable verbose logging")
    parser.add_argument("--detect-only", "-d", action="store_true",
                        help="Only detect document types without full processing")
                        
    args = parser.parse_args()
    
    # Set verbose logging if requested
    if args.verbose:
        logger.setLevel(logging.DEBUG)
        for handler in logger.handlers:
            handler.setLevel(logging.DEBUG)
        logger.debug("Verbose logging enabled")

    # Get batch size from DB or command line
    db_manager = DBManager()
    default_limit = 500
    
    # Command line overrides DB setting
    if args.limit is not None:
        batch_size = args.limit
        logger.info(f"Using command line batch size: {batch_size}")
    else:
        batch_size = get_batch_size_from_settings(db_manager, default_limit)
    
    # Normalize format name
    output_format = args.format.lower()
    if output_format == "md":
        output_format = "markdown"
        
    logger.info(f"Starting document loader with limit={batch_size}, format={output_format}, retry_failed={args.retry}")
    
    if args.detect_only:
        # If implemented, document type detection mode
        logger.info("Document type detection only mode")
        # This would be a separate function that doesn't do full extraction
    else:
        # Run the loader
        processed_count = load_documents(limit=batch_size, output_format=output_format, retry_failed=args.retry)
        logger.info(f"Document loading completed. Processed {processed_count} documents.")

if __name__ == "__main__":
    main()
