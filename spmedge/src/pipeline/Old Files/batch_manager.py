#!/usr/bin/env python3
"""
Batch Manager - Manages document batches by type and initializes processing pipeline.
- Processes documents from unprocessed directory
- Requires document type specification
- Standardizes file naming and organization
- Prepares batches for the document loader stage
- Provides pipeline management and reset capabilities
"""

import os
import sys
import json
import logging
import argparse
import shutil
import re
from pathlib import Path
from datetime import datetime
from typing import Optional, List, Dict, Any, Tuple

# Fix the import path issue - use absolute imports
sys.path.insert(0, '/Users/toddlebaron/prizym/spmedge')

# Now import using the absolute path
from config.config import config
from src.pipeline.db_integration import DBManager
from src.pipeline.pipeline_processor import PipelineProcessor, PipelineStage

# Ensure logs directory exists
LOG_DIR = config.LOG_DIR
os.makedirs(LOG_DIR, exist_ok=True)
LOG_FILE = os.path.join(LOG_DIR, "batch_manager.log")

# Configure logging
logger = logging.getLogger("batch_manager")
logger.propagate = False  # Prevent propagation to parent loggers

if not logger.handlers:
    logger.setLevel(logging.INFO)
    formatter = logging.Formatter("%(asctime)s - %(levelname)s - %(message)s")
    
    # Console handler
    console_handler = logging.StreamHandler()
    console_handler.setFormatter(formatter)
    logger.addHandler(console_handler)
    
    # File handler
    file_handler = logging.FileHandler(LOG_FILE, encoding="utf-8")
    file_handler.setFormatter(formatter)
    logger.addHandler(file_handler)

class BatchManager:
    """Manages document batches by type and initializes processing pipeline."""
    
    def __init__(self):
        """Initialize BatchManager with database connection and processor."""
        self.db_manager = DBManager()
        self.processor = PipelineProcessor(PipelineStage.INPUT)
    
    def sanitize_filename(self, filename: str) -> str:
        """Clean and standardize filenames for processing."""
        # Get base name and extension
        base_name = Path(filename).stem
        extension = Path(filename).suffix
        
        # Remove special characters, replace spaces with underscores
        clean_name = re.sub(r'[^\w\-\.]', '_', base_name)
        # Remove multiple underscores
        clean_name = re.sub(r'_+', '_', clean_name)
        # Limit length
        clean_name = clean_name[:100]
        
        return f"{clean_name}{extension}"
    
    def verify_document_type(self, doc_type: str) -> int:
        """Verify document type exists in database and return its ID."""
        try:
            # Check if document type exists
            self.db_manager.cursor.execute("SELECT id FROM document_types WHERE name = %s;", (doc_type,))
            result = self.db_manager.cursor.fetchone()
            
            if not result:
                logger.error(f"❌ Invalid document type: {doc_type}")
                raise ValueError(f"Document type '{doc_type}' does not exist in the database")
            
            return result[0]  # Return document type ID
            
        except Exception as e:
            logger.error(f"❌ Error verifying document type: {e}")
            raise
    def process_single_document(self, file_path: Path, doc_type: str, doc_type_id: int, 
                               batch_id: str, stage_input_dir: Path, archive_dir: Path, 
                               archive: bool = False, delete: bool = False) -> Optional[Dict[str, Any]]:
        """Process a single document and return its metadata if successful."""
        try:
            # Clean and standardize filename
            clean_filename = self.sanitize_filename(file_path.name)
            
            # Generate new filename with batch info and .json extension
            base_name = Path(clean_filename).stem
            new_filename = f"{doc_type}_{base_name}.json"
            
            # Read and convert file content to JSON structure
            file_content = self.extract_file_content(file_path)
            
            # Create JSON structure with metadata
            json_content = {
                "filename": file_path.name,
                "file_type": file_path.suffix.lower(),
                "file_size": file_path.stat().st_size,
                "extraction_time": datetime.now().isoformat(),
                "content": file_content
            }
            
            # Insert document record in database
            self.db_manager.cursor.execute(
                "INSERT INTO documents (id, name, original_filename, document_type_id, batch_id) VALUES (gen_random_uuid(), %s, %s, %s, %s) RETURNING id;",
                (new_filename, file_path.name, doc_type_id, batch_id)
            )
            document_id = self.db_manager.cursor.fetchone()[0]
            
            # CRITICAL: Force commit and create a new transaction before updating pipeline
            self.db_manager.conn.commit()
            
            logger.info(f"✅ Registered document {document_id}: {new_filename} (originally {file_path.name})")
            
            # Save JSON to stage_input directory
            stage_file_path = stage_input_dir / new_filename
            with open(stage_file_path, 'w', encoding='utf-8') as f:
                json.dump(json_content, f, indent=2)
            
            logger.info(f"📁 Created JSON file at stage directory: {stage_file_path}")
            
            # Handle original file based on options
            if archive:
                archive_path = archive_dir / file_path.name
                shutil.copy2(str(file_path), str(archive_path))
                logger.info(f"📦 Archived original file to: {archive_path}")
    
            if delete or archive:  # If archived, we can also delete
                os.unlink(str(file_path))
                logger.info(f"🗑️ Deleted original file: {file_path}")
    
            # Return document metadata
            return {
                "id": document_id,
                "name": new_filename,
                "original_filename": file_path.name,
                "file_path": str(stage_file_path),
                "original_path": str(file_path),
                "document_type_id": doc_type_id,
                "document_type": doc_type,
                "status": "completed",
                "created_at": datetime.now().isoformat()
            }
    
        except Exception as e:
            logger.error(f"❌ Failed to process {file_path.name}: {e}")
            return None
    
    def extract_file_content(self, file_path: Path) -> str:
        """Extract content from a file based on its type."""
        file_type = file_path.suffix.lower()
        
        try:
            # Text-based files
            if file_type in ['.txt', '.md', '.csv', '.json']:
                with open(file_path, 'r', encoding='utf-8') as f:
                    return f.read()
                    
            # PDF files
            elif file_type == '.pdf':
                try:
                    from PyPDF2 import PdfReader
                    content = ""
                    with open(file_path, 'rb') as f:
                        reader = PdfReader(f)
                        for page in reader.pages:
                            content += page.extract_text() + "\n"
                    return content
                except ImportError:
                    logger.warning("PyPDF2 not installed. Cannot extract PDF content.")
                    return f"[PDF file: {file_path.name}]"
                    
            # DOCX files
            elif file_type == '.docx':
                try:
                    import docx
                    doc = docx.Document(file_path)
                    return "\n".join([para.text for para in doc.paragraphs])
                except ImportError:
                    logger.warning("python-docx not installed. Cannot extract DOCX content.")
                    return f"[DOCX file: {file_path.name}]"
                    
            # XLSX files
            elif file_type == '.xlsx':
                try:
                    import pandas as pd
                    excel_file = pd.ExcelFile(file_path)
                    sheets = excel_file.sheet_names
                    
                    content = []
                    for sheet in sheets:
                        df = pd.read_excel(excel_file, sheet_name=sheet)
                        content.append(f"Sheet: {sheet}")
                        content.append(df.to_string(index=False))
                        content.append("\n")
                        
                    return "\n".join(content)
                except ImportError:
                    logger.warning("pandas not installed. Cannot extract XLSX content.")
                    return f"[XLSX file: {file_path.name}]"
                    
            # PPTX files
            elif file_type == '.pptx':
                try:
                    from pptx import Presentation
                    text_content = []
                    
                    presentation = Presentation(file_path)
                    for i, slide in enumerate(presentation.slides):
                        slide_text = []
                        slide_text.append(f"Slide {i+1}")
                        
                        if slide.shapes.title and slide.shapes.title.text:
                            slide_text.append(f"Title: {slide.shapes.title.text}")
                        
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text:
                                slide_text.append(shape.text)
                        
                        text_content.append("\n".join(slide_text))
                    
                    return "\n\n".join(text_content)
                except ImportError:
                    logger.warning("python-pptx not installed. Cannot extract PPTX content.")
                    return f"[PowerPoint file: {file_path.name}]"
            
            # Binary or unsupported files
            else:
                return f"[Binary file: {file_path.name}]"
                
        except UnicodeDecodeError:
            logger.warning(f"Unicode decode error with {file_path}")
            return f"[Binary file: {file_path.name}]"
        except Exception as e:
            logger.error(f"Error extracting content from {file_path}: {e}")
            return f"[Error extracting content: {str(e)}]"
                                
    def update_pipeline_entries(self, processed_documents: List[Dict[str, Any]], batch_id: str, document_type_id: int):
        """
        Update pipeline entries for multiple documents in a separate transaction.
        
        Args:
            processed_documents: List of document dictionaries
            batch_id: The batch ID
            document_type_id: The document type ID
        
        Returns:
            int: Number of successful updates
        """
        success_count = 0
        
        # Create a separate cursor for this operation
        connection = self.db_manager.conn
        previous_autocommit = connection.autocommit
        connection.autocommit = True
        
        try:
            cursor = connection.cursor()
            
            for doc in processed_documents:
                try:
                    document_id = doc["id"]
                    # Insert pipeline entry
                    cursor.execute(
                        """
                        INSERT INTO processing_pipeline 
                        (document_id, document_type_id, pipeline_stage, status, batch_id, updated_at)
                        VALUES (%s, %s, %s, %s, %s, NOW())
                        ON CONFLICT (document_id, pipeline_stage) 
                        DO UPDATE SET status = EXCLUDED.status, updated_at = NOW();
                        """,
                        (document_id, document_type_id, "input", "completed", batch_id)
                    )
                    success_count += 1
                except Exception as e:
                    logger.error(f"❌ Failed to update pipeline for document {doc.get('id')}: {e}")
                    # Continue with next document
            
            logger.info(f"✅ Updated pipeline entries for {success_count} of {len(processed_documents)} documents")
            return success_count
            
        finally:
            # Restore autocommit mode
            connection.autocommit = previous_autocommit
    
    def process_document_batch(self, doc_type: str, archive: bool = False, delete: bool = False) -> Optional[str]:
        """
        Process a batch of documents of the specified type.

        Args:
            doc_type: Document type for this batch
            archive: Whether to archive original files
            delete: Whether to delete original files

        Returns:
            str: Batch ID if successful, None otherwise
        """
        if not doc_type:
            logger.error("❌ Document type must be specified!")
            return None

        batch_id = None

        try:
            # Verify document type exists in database
            document_type_id = self.verify_document_type(doc_type)

            # Get directories
            dirs = self.processor.get_base_dirs()
            unprocessed_dir = dirs["unprocessed"]
            stage_input_dir = dirs["stage_input"]
            archive_dir = dirs["archive"]

            # Ensure directories exist
            unprocessed_dir.mkdir(parents=True, exist_ok=True)
            stage_input_dir.mkdir(parents=True, exist_ok=True)
            archive_dir.mkdir(parents=True, exist_ok=True)

            # Find files in unprocessed directory
            files = list(unprocessed_dir.glob("*.*"))

            if not files:
                logger.info("ℹ️ No documents found in unprocessed directory.")
                return None

            logger.info(f"🔍 Found {len(files)} documents in: {unprocessed_dir}")

            # Create batch with auto-commit to ensure it's available
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            batch_name = f"batch_{doc_type}_{timestamp}"
            batch_id = self.processor.record_batch_processing(
                batch_name=batch_name,
                document_count=len(files),
                status="processing"
            )
            
            processed_documents = []
            success_count = 0

            # Process each document individually (each in its own transaction)
            for file_path in files:
                doc_metadata = self.process_single_document(
                    file_path=file_path,
                    doc_type=doc_type,
                    doc_type_id=document_type_id,
                    batch_id=batch_id,
                    stage_input_dir=stage_input_dir,
                    archive_dir=archive_dir,
                    archive=archive,
                    delete=delete
                )
                
                if doc_metadata:
                    processed_documents.append(doc_metadata)
                    success_count += 1

            # If we have processed documents, save batch metadata
            if processed_documents:
                self.processor.save_document_batch(documents=processed_documents, batch_name=batch_name)
                logger.info(f"🎉 Successfully processed {success_count} of {len(files)} documents of type '{doc_type}'")

                # Now do a separate pass to update pipeline entries
                # This prevents foreign key constraint issues by doing it after all documents are committed
                pipeline_updates = self.update_pipeline_entries(
                    processed_documents=processed_documents,
                    batch_id=batch_id,
                    document_type_id=document_type_id
                )
                
                # Finalize batch status
                if success_count == len(files):
                    self.processor.finalize_batch(batch_id, "completed")
                elif success_count > 0:
                    self.processor.finalize_batch(batch_id, "partial")
                else:
                    self.processor.finalize_batch(batch_id, "failed")
                    batch_id = None
            else:
                logger.warning("⚠️ No documents were successfully processed.")
                batch_id = None
                    
            logger.info("🏁 Batch processing complete")
            return batch_id

        except Exception as e:
            logger.error(f"❌ Batch processing failed: {e}")
            return None
    
    def reset_pipeline(self, stage: str = None, batch_id: str = None):
        """
        Reset the pipeline - either a specific stage or the entire pipeline.
        
        Args:
            stage: Specific stage to reset (None for all stages)
            batch_id: Optional batch ID to limit the reset scope
        
        Returns:
            bool: Success or failure
        """
        try:
            if batch_id:
                logger.info(f"Resetting pipeline for batch: {batch_id}")
            
            if stage:
                # Reset only the specified stage
                result = self.processor.reset_pipeline_stage(stage, batch_id)
                logger.info(f"✅ Reset pipeline stage: {stage}" + (f" for batch {batch_id}" if batch_id else ""))
            else:
                # Reset all stages in reverse order (to prevent foreign key issues)
                for s in ["report", "process", "clean", "load", "input"]:
                    self.processor.reset_pipeline_stage(s, batch_id)
                logger.info("✅ Reset entire pipeline" + (f" for batch {batch_id}" if batch_id else ""))
            
            if batch_id:
                # Keep track of autocommit state
                previous_autocommit = self.db_manager.conn.autocommit
                self.db_manager.conn.autocommit = False
                
                try:
                    # Delete the batch if no documents reference it
                    self.db_manager.cursor.execute(
                        "SELECT COUNT(*) FROM documents WHERE batch_id = %s;", 
                        (batch_id,)
                    )
                    doc_count = self.db_manager.cursor.fetchone()[0]
                    
                    if doc_count == 0:
                        # Safe to delete the batch
                        self.db_manager.cursor.execute(
                            "DELETE FROM processing_batches WHERE batch_id = %s;", 
                            (batch_id,)
                        )
                        self.db_manager.conn.commit()
                        logger.info(f"✅ Deleted batch: {batch_id}")
                    else:
                        logger.warning(f"⚠️ Batch {batch_id} still has {doc_count} documents, not deleting batch record")
                except Exception as e:
                    logger.error(f"❌ Failed to delete batch: {e}")
                    self.db_manager.conn.rollback()
                finally:
                    # Restore autocommit mode
                    self.db_manager.conn.autocommit = previous_autocommit
                    
            return True
        except Exception as e:
            logger.error(f"❌ Failed to reset pipeline: {e}")
            return False

    def list_active_batches(self):
        """List all active batches in the pipeline."""
        try:
            self.db_manager.cursor.execute("""
                SELECT b.batch_id, b.batch_name, b.document_count, b.status, b.created_at, b.completed_at,
                      (SELECT COUNT(*) FROM documents d WHERE d.batch_id = b.batch_id) as doc_count
                FROM processing_batches b
                WHERE b.status IN ('processing', 'partial')
                ORDER BY b.created_at DESC;
            """)
            
            batches = self.db_manager.cursor.fetchall()
            
            if not batches:
                logger.info("No active batches found")
                return
                
            logger.info("Active batches:")
            logger.info(f"{'ID':<36} | {'Batch Name':<30} | {'Status':<12} | {'Doc Count':<9} | {'Created At'}")
            logger.info("-" * 100)
            
            for batch in batches:
                batch_id = batch[0]
                batch_name = batch[1]
                status = batch[3]
                doc_count = batch[6]
                created_at = batch[4].strftime("%Y-%m-%d %H:%M") if batch[4] else "Unknown"
                
                logger.info(f"{batch_id:<36} | {batch_name:<30} | {status:<12} | {doc_count:<9} | {created_at}")
                
            # Also list counts by pipeline stage
            logger.info("\nPipeline stage document counts:")
            self.db_manager.cursor.execute("""
                SELECT pipeline_stage, status, COUNT(*)
                FROM processing_pipeline
                GROUP BY pipeline_stage, status
                ORDER BY pipeline_stage, status;
            """)
            
            stage_counts = self.db_manager.cursor.fetchall()
            
            if stage_counts:
                logger.info(f"{'Stage':<10} | {'Status':<12} | {'Count'}")
                logger.info("-" * 40)
                
                for stage, status, count in stage_counts:
                    logger.info(f"{stage:<10} | {status:<12} | {count}")
            
        except Exception as e:
            logger.error(f"❌ Failed to list active batches: {e}")

    def cleanup_orphans(self) -> Tuple[int, int]:
        """Clean up orphaned documents and records."""
        try:
            # Keep track of autocommit state
            previous_autocommit = self.db_manager.conn.autocommit
            self.db_manager.conn.autocommit = False
            
            try:
                # First, clean up documents with no pipeline entries
                doc_count = self.processor.cleanup_orphaned_documents()
                
                # Clean up empty batches
                self.db_manager.cursor.execute("""
                    DELETE FROM processing_batches pb
                    WHERE NOT EXISTS (
                        SELECT 1 FROM documents d
                        WHERE d.batch_id = pb.batch_id
                    )
                    RETURNING batch_id, batch_name;
                """)
                
                deleted_batches = self.db_manager.cursor.fetchall()
                batch_count = len(deleted_batches)
                
                if batch_count > 0:
                    logger.info(f"✅ Cleaned up {batch_count} orphaned batches")
                    for batch_id, name in deleted_batches:
                        logger.debug(f"Deleted orphaned batch: {batch_id} ({name})")
                
                self.db_manager.conn.commit()
                return doc_count, batch_count
            except Exception as e:
                logger.error(f"❌ Failed to clean up orphans: {e}")
                self.db_manager.conn.rollback()
                return 0, 0
            finally:
                # Restore autocommit mode
                self.db_manager.conn.autocommit = previous_autocommit
                
        except Exception as e:
            logger.error(f"❌ Failed to clean up orphans: {e}")
            return 0, 0

    def batch_status(self, batch_id: str):
        """Display detailed status for a specific batch."""
        status = self.processor.get_batch_status(batch_id)
        
        if not status:
            logger.error(f"❌ Batch {batch_id} not found or error occurred")
            return
            
        logger.info(f"Batch Status: {status['batch_name']} ({status['batch_id']})")
        logger.info(f"Status: {status['status']}")
        logger.info(f"Documents: {status['document_count']}")
        logger.info(f"Created: {status['created_at']}")
        
        if status['completed_at']:
            logger.info(f"Completed: {status['completed_at']}")
        
        logger.info("\nPipeline Stage Status:")
        
        if not status['pipeline_stages']:
            logger.info("No pipeline stage data found")
        else:
            for stage, stage_data in status['pipeline_stages'].items():
                stage_str = f"{stage}: "
                for status_name, count in stage_data.items():
                    stage_str += f"{status_name}={count} "
                logger.info(stage_str)

def main():
    """Main function to handle batch document processing."""
    batch_manager = BatchManager()
    
    parser = argparse.ArgumentParser(description="Batch Manager for Document Processing Pipeline")
    
    # Create subparsers for different commands
    subparsers = parser.add_subparsers(dest="command", help="Command to execute")
    
    # Process command
    process_parser = subparsers.add_parser("process", help="Process a batch of documents")
    process_parser.add_argument("doc_type", help="Document type to process (REQUIRED)")
    process_parser.add_argument("--archive", action="store_true", help="Archive original files after processing")
    process_parser.add_argument("--delete", action="store_true", help="Delete original files after processing")
    
    # Reset command
    reset_parser = subparsers.add_parser("reset", help="Reset pipeline stages")
    reset_parser.add_argument("--stage", choices=["input", "load", "clean", "process", "report"], 
                              help="Specific stage to reset (omit for all stages)")
    reset_parser.add_argument("--batch", help="Specific batch ID to reset (optional)")
    
    # List command
    list_parser = subparsers.add_parser("list", help="List active batches")
    
    # Cleanup command
    cleanup_parser = subparsers.add_parser("cleanup", help="Clean up orphaned documents and batches")
    
    # Status command
    status_parser = subparsers.add_parser("status", help="Check status of a specific batch")
    status_parser.add_argument("batch_id", help="ID of the batch to check")
    
    args = parser.parse_args()
    
    if not args.command:
        parser.print_help()
        return
        
    if args.command == "process":
        if not args.doc_type:
            process_parser.error("Document type is required")
        batch_manager.process_document_batch(args.doc_type, args.archive, args.delete)
    
    elif args.command == "reset":
        batch_manager.reset_pipeline(args.stage, args.batch)
    
    elif args.command == "list":
        batch_manager.list_active_batches()
    
    elif args.command == "cleanup":
        doc_count, batch_count = batch_manager.cleanup_orphans()
        logger.info(f"Cleanup completed: {doc_count} documents and {batch_count} batches removed")
    
    elif args.command == "status":
        batch_manager.batch_status(args.batch_id)

if __name__ == "__main__":
    main()